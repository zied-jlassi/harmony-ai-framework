"""
PowerPoint Parser - 100% Local RGPD Compliant
Uses python-pptx - NO cloud APIs

OPTIMIZATIONS:
- Slide selection: --slides "1-10" or "1,3,5-8"
"""

from pptx import Presentation
import json
import argparse
from typing import Optional, List
from pathlib import Path


class PPTXParser:
    """Parse PowerPoint 100% local - RGPD compliant"""

    def __init__(self):
        pass

    def parse(self, file_path: str, slides: Optional[str] = None) -> dict:
        """
        Extract text, structure and notes from PowerPoint.

        Args:
            file_path: Path to PPTX file
            slides: Slide selection (e.g., "5", "1-10", "1,3,5-8")

        Returns:
            dict with metadata and slides content
        """
        try:
            prs = Presentation(file_path)
        except Exception as e:
            raise ValueError(f"Invalid PPTX file: {file_path} - {str(e)}")

        total_slides = len(prs.slides)

        # Parse slide selection
        slide_indices = self._parse_slide_range(slides, total_slides)

        result = {
            "source": str(Path(file_path).name),
            "parser": "python-pptx",
            "rgpd_compliant": True,
            "metadata": {
                "total_slides": total_slides,
                "parsed_slides": len(slide_indices),
                "slide_range": slides if slides else "all"
            },
            "slides": []
        }

        for idx in slide_indices:
            slide = prs.slides[idx]
            slide_data = self._extract_slide(slide, idx + 1)
            result["slides"].append(slide_data)

        return result

    def _parse_slide_range(self, slides: Optional[str], total: int) -> List[int]:
        """
        Parse slide specification.

        Args:
            slides: Slide specification:
                - None: all slides
                - "5": slide 5 only
                - "1-10": slides 1 to 10
                - "1,3,5-8": slides 1, 3, and 5 to 8
            total: Total number of slides

        Returns:
            List of 0-indexed slide indices
        """
        if slides is None:
            return list(range(total))

        indices = set()
        for part in slides.split(','):
            part = part.strip()
            if '-' in part:
                start, end = part.split('-', 1)
                start = max(1, int(start))
                end = min(total, int(end))
                indices.update(range(start - 1, end))
            else:
                slide = int(part)
                if 1 <= slide <= total:
                    indices.add(slide - 1)

        return sorted(indices)

    def _extract_slide(self, slide, slide_num: int) -> dict:
        """Extract content from a single slide."""
        slide_data = {
            "number": slide_num,
            "title": None,
            "content": [],
            "notes": None,
            "shapes_count": len(slide.shapes)
        }

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape.is_placeholder:
                    try:
                        if shape.placeholder_format.type == 1:
                            slide_data["title"] = shape.text.strip()
                        else:
                            slide_data["content"].append(shape.text.strip())
                    except:
                        slide_data["content"].append(shape.text.strip())
                else:
                    slide_data["content"].append(shape.text.strip())

            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                slide_data["content"].append({"type": "table", "data": table_data})

        if slide.has_notes_slide:
            try:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame and notes_frame.text.strip():
                    slide_data["notes"] = notes_frame.text.strip()
            except:
                pass

        return slide_data

    def parse_to_json(self, file_path: str, slides: Optional[str] = None, output_path: Optional[str] = None) -> str:
        """
        Parse PPTX and save as JSON.

        Args:
            file_path: Path to PPTX file
            slides: Slide selection
            output_path: Optional output path for JSON

        Returns:
            JSON string
        """
        result = self.parse(file_path, slides)
        json_str = json.dumps(result, ensure_ascii=False, indent=2)

        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(json_str)

        return json_str


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="PowerPoint Parser - 100% Local RGPD Compliant",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python pptx_parser.py presentation.pptx                 # All slides
  python pptx_parser.py presentation.pptx --slides 5      # Slide 5 only
  python pptx_parser.py presentation.pptx --slides 1-10   # Slides 1 to 10
  python pptx_parser.py presentation.pptx --slides 1,3,5-8  # Slides 1, 3, 5-8
  python pptx_parser.py presentation.pptx -o output.json  # Save to file
        """
    )
    parser.add_argument("file", help="Path to PPTX file")
    parser.add_argument("--slides", "-s", help="Slide selection (e.g., '1-10', '1,3,5-8')")
    parser.add_argument("--output", "-o", help="Output JSON file path")

    args = parser.parse_args()

    pptx_parser = PPTXParser()
    result = pptx_parser.parse_to_json(args.file, args.slides, args.output)
    print(result)
